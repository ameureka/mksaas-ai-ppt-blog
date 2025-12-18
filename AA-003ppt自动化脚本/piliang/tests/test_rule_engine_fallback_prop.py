from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation

from factory.workshops.workshopC import WorkshopC, WorkshopCConfig
from factory.types import EtlOutput


def _create_test_pptx(path: Path) -> None:
	prs = Presentation()
	slide = prs.slides.add_slide(prs.slide_layouts[6])
	prs.save(path)


def _setup_project_root(tmp_path: Path, mapping_content: str) -> Path:
	"""Setup a minimal project root with required config files."""
	project_root = tmp_path / 'project'
	project_root.mkdir()
	
	configs_dir = project_root / 'configs'
	configs_dir.mkdir()
	(configs_dir / 'category-mapping.yaml').write_text(mapping_content, encoding='utf-8')
	(configs_dir / 'forbidden-keywords.txt').write_text('', encoding='utf-8')
	
	templates_dir = project_root / 'templates'
	templates_dir.mkdir()
	(templates_dir / 'ai_prompt_template.md').write_text(
		'# Context\nTitle: {{title}}\nTags: {{original_tags}}\n# Text\n{{extracted_text}}\n',
		encoding='utf-8',
	)
	
	return project_root


def test_property12_rule_priority_vs_ai(tmp_path: Path, monkeypatch) -> None:
	# Rules cover "报告" -> report; missing cases trigger AI fallback.
	project_root = _setup_project_root(
		tmp_path,
		"categories:\n  report:\n    priority: 10\n    keywords:\n      - 报告\n",
	)

	def fake_run(args, capture_output, text, timeout=None):  # noqa: ANN001
		output_path = Path(args[-1])
		out = {
			'ai_summary': 'This is a test summary for the presentation',
			'ai_keywords': ['k1', 'k2'],
			'ai_scenario': 'business presentation',
			'ai_color_scheme': 'blue',
			'ai_structure_features': 'sections',
			'ai_template_features': 'clean',
			'ppthub_category': 'general',
			'language': '中文',
		}
		output_path.parent.mkdir(parents=True, exist_ok=True)
		output_path.write_text(json.dumps(out, ensure_ascii=False), encoding='utf-8')

		class R:
			returncode = 0
			stdout = ''
			stderr = ''
		return R()

	import subprocess  # noqa: WPS433
	monkeypatch.setattr(subprocess, 'run', fake_run)

	config = WorkshopCConfig(
		project_root=project_root,
		prompt_output_dir=tmp_path / 'prompts',
		ai_output_dir=tmp_path / 'ai',
		ai_command=['cli', '{prompt}', '{output}'],
	)
	wc = WorkshopC(config)

	# Create test PPTX
	pptx_path = tmp_path / 'test.pptx'
	_create_test_pptx(pptx_path)

	etl_rule_hit = EtlOutput(
		aid='hit',
		channel_id='chan',
		local_pptx_path=pptx_path,
		local_cover_path=None,
		pages_count=1,
		file_size_kb=1,
		meta={'title': '年度报告', 'original_tags': [], 'channel_name': 'test'},
	)
	etl_rule_miss = EtlOutput(
		aid='miss',
		channel_id='chan',
		local_pptx_path=pptx_path,
		local_cover_path=None,
		pages_count=1,
		file_size_kb=1,
		meta={'title': '未知标题', 'original_tags': [], 'channel_name': 'test'},
	)

	res_hit = wc.run(etl_rule_hit)
	assert res_hit['ppthub_category'] == 'report'
	assert res_hit['ai_meta'] is not None  # AI 始终调用以获取丰富内容

	res_miss = wc.run(etl_rule_miss)
	assert res_miss['ppthub_category'] == 'general'
	assert res_miss['ai_meta'] is not None
