"""
TextExtractor 模块单元测试
"""

from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from factory.ai.text_extractor import TextExtractor


@pytest.fixture
def sample_pptx(tmp_path: Path) -> Path:
	"""创建测试用的 PPTX 文件"""
	prs = Presentation()
	prs.slide_width = Inches(10)
	prs.slide_height = Inches(7.5)

	# 第一页：中文内容
	slide1 = prs.slides.add_slide(prs.slide_layouts[5])
	textbox1 = slide1.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
	textbox1.text = '这是一个测试标题'
	textbox2 = slide1.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
	textbox2.text = '这是测试内容，包含中文字符。'

	# 第二页：英文内容
	slide2 = prs.slides.add_slide(prs.slide_layouts[5])
	textbox3 = slide2.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
	textbox3.text = 'This is a test title'
	textbox4 = slide2.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
	textbox4.text = 'This is test content with English text.'

	# 第三页：包含占位符
	slide3 = prs.slides.add_slide(prs.slide_layouts[5])
	textbox5 = slide3.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
	textbox5.text = 'Click to add title'
	textbox6 = slide3.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
	textbox6.text = 'Lorem ipsum dolor sit amet'

	pptx_path = tmp_path / 'test.pptx'
	prs.save(str(pptx_path))
	return pptx_path


def test_extract_success(sample_pptx: Path) -> None:
	"""测试成功提取文本"""
	extractor = TextExtractor()
	result = extractor.extract(sample_pptx)

	assert result.full_text
	assert result.char_count > 0
	assert len(result.slide_texts) > 0
	assert '这是一个测试标题' in result.full_text
	assert 'This is a test title' in result.full_text


def test_extract_file_not_found() -> None:
	"""测试文件不存在"""
	extractor = TextExtractor()
	with pytest.raises(FileNotFoundError):
		extractor.extract(Path('/nonexistent/file.pptx'))


def test_extract_invalid_file(tmp_path: Path) -> None:
	"""测试无效的 PPTX 文件"""
	invalid_file = tmp_path / 'invalid.pptx'
	invalid_file.write_text('not a valid pptx', encoding='utf-8')

	extractor = TextExtractor()
	with pytest.raises(ValueError, match='Failed to parse'):
		extractor.extract(invalid_file)


def test_chinese_ratio_calculation(sample_pptx: Path) -> None:
	"""测试中文字符占比计算"""
	extractor = TextExtractor()
	result = extractor.extract(sample_pptx)

	assert 0.0 <= result.chinese_ratio <= 1.0
	assert result.chinese_ratio > 0  # 包含中文内容


def test_english_ratio_calculation(sample_pptx: Path) -> None:
	"""测试英文字符占比计算"""
	extractor = TextExtractor()
	result = extractor.extract(sample_pptx)

	assert 0.0 <= result.english_ratio <= 1.0
	assert result.english_ratio > 0  # 包含英文内容


def test_filter_placeholders() -> None:
	"""测试占位符过滤"""
	extractor = TextExtractor()

	# 测试各种占位符
	test_cases = [
		('Click to add title', ''),
		('Lorem ipsum dolor', ''),
		('CLICK TO EDIT', ''),
		('Add your text here', ''),
		('Type here to begin', ''),
		('正常文本\nClick to add\n更多文本', '正常文本\n更多文本'),
		('这是正常内容', '这是正常内容'),
	]

	for input_text, expected in test_cases:
		result = extractor.filter_placeholders(input_text)
		assert result == expected, f'Failed for input: {input_text}'


def test_filter_placeholders_case_insensitive() -> None:
	"""测试占位符过滤不区分大小写"""
	extractor = TextExtractor()

	test_cases = [
		'lorem ipsum',
		'Lorem Ipsum',
		'LOREM IPSUM',
		'LoReM iPsUm',
	]

	for text in test_cases:
		result = extractor.filter_placeholders(text)
		assert result == '', f'Failed to filter: {text}'


def test_extract_empty_pptx(tmp_path: Path) -> None:
	"""测试空 PPTX 文件"""
	prs = Presentation()
	prs.slides.add_slide(prs.slide_layouts[5])  # 空白页

	pptx_path = tmp_path / 'empty.pptx'
	prs.save(str(pptx_path))

	extractor = TextExtractor()
	result = extractor.extract(pptx_path)

	assert result.full_text == ''
	assert result.char_count == 0
	assert result.chinese_ratio == 0.0
	assert result.english_ratio == 0.0


def test_extract_chinese_only(tmp_path: Path) -> None:
	"""测试纯中文内容"""
	prs = Presentation()
	slide = prs.slides.add_slide(prs.slide_layouts[5])
	textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
	textbox.text = '这是纯中文内容测试'

	pptx_path = tmp_path / 'chinese.pptx'
	prs.save(str(pptx_path))

	extractor = TextExtractor()
	result = extractor.extract(pptx_path)

	assert result.chinese_ratio > 0.5
	assert '这是纯中文内容测试' in result.full_text


def test_extract_english_only(tmp_path: Path) -> None:
	"""测试纯英文内容"""
	prs = Presentation()
	slide = prs.slides.add_slide(prs.slide_layouts[5])
	textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
	textbox.text = 'This is pure English content test'

	pptx_path = tmp_path / 'english.pptx'
	prs.save(str(pptx_path))

	extractor = TextExtractor()
	result = extractor.extract(pptx_path)

	assert result.english_ratio > 0.5
	assert 'This is pure English content test' in result.full_text
