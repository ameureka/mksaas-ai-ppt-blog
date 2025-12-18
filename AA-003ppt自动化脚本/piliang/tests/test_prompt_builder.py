from __future__ import annotations

import json
from pathlib import Path

from factory.ai.prompt_builder import PromptBuilder, PromptBuilderConfig


def test_prompt_builder_extracts_text(tmp_path: Path) -> None:
	from pptx import Presentation

	pptx_path = tmp_path / 'sample.pptx'
	prs = Presentation()
	layout = prs.slide_layouts[5]
	s1 = prs.slides.add_slide(layout)
	s1.shapes.title.text = 'Hello World'
	s2 = prs.slides.add_slide(layout)
	s2.shapes.title.text = 'Slide 2 content'
	prs.save(pptx_path)

	config = PromptBuilderConfig(output_dir=tmp_path / 'prompts')
	builder = PromptBuilder(config)
	meta = {'aid': '1', 'title': 'Demo', 'channel_id': 'chan'}
	out_path = builder.build(aid='1', title='Demo', meta=meta, pptx_path=pptx_path)

	assert out_path.exists()
	data = json.loads(out_path.read_text(encoding='utf-8'))
	assert data['aid'] == '1'
	assert 'Hello World' in data['extracted_text']
	assert 'Slide 2 content' in data['extracted_text']
