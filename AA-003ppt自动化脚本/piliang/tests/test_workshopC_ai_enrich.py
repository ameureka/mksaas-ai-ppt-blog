from __future__ import annotations

import json
import sqlite3
from pathlib import Path

import pytest
from pptx import Presentation

from factory.db.init_db import init_db
from factory.db.dao import get_stage_status
from factory.stages import StageName, StageStatus
from factory.types import EtlOutput
from factory.workshops.workshopC import WorkshopC, WorkshopCConfig, run_ai_enrich


def _create_test_pptx(path: Path) -> None:
	prs = Presentation()
	slide = prs.slides.add_slide(prs.slide_layouts[6])
	prs.save(path)


def _setup_project_root(tmp_path: Path, mapping_content: str) -> Path:
	"""Setup a minimal project root with required config files."""
	project_root = tmp_path / 'project'
	project_root.mkdir()
	
	# Create configs dir
	configs_dir = project_root / 'configs'
	configs_dir.mkdir()
	(configs_dir / 'category-mapping.yaml').write_text(mapping_content, encoding='utf-8')
	(configs_dir / 'forbidden-keywords.txt').write_text('', encoding='utf-8')
	
	# Create templates dir
	templates_dir = project_root / 'templates'
	templates_dir.mkdir()
	(templates_dir / 'ai_prompt_template.md').write_text(
		'# Context\nTitle: {{title}}\nTags: {{original_tags}}\n# Text\n{{extracted_text}}\n',
		encoding='utf-8',
	)
	
	return project_root


def test_rule_engine_wins_without_ai(tmp_path: Path) -> None:
	project_root = _setup_project_root(
		tmp_path,
		"categories:\n  report:\n    priority: 10\n    keywords:\n      - 报告\n      - 年度\n",
	)
	config = WorkshopCConfig(
		project_root=project_root,
		prompt_output_dir=tmp_path / 'prompts',
		ai_output_dir=tmp_path / 'ai',
		ai_command=['echo', '{prompt}', '{output}'],
	)
	wc = WorkshopC(config)
	
	# Create test PPTX
	pptx_path = tmp_path / 'test.pptx'
	_create_test_pptx(pptx_path)
	
	etl_out = EtlOutput(
		aid='a1',
		channel_id='chan',
		local_pptx_path=pptx_path,
		local_cover_path=None,
		pages_count=1,
		file_size_kb=1,
		meta={'title': '年度报告', 'original_tags': ['报告'], 'channel_name': 'test'},
	)
	db = tmp_path / 'assets.db'
	init_db(db)
	conn = sqlite3.connect(db)
	conn.row_factory = sqlite3.Row
	try:
		# 预先创建 AI 输出文件，因为 echo 命令不会创建它
		ai_output_path = tmp_path / 'ai' / 'a1.json'
		ai_output_path.parent.mkdir(parents=True, exist_ok=True)
		fake_ai_output = {
			'ai_summary': 'Test summary',
			'ai_keywords': ['test', 'keyword'],
			'ai_scenario': 'test scenario',
			'ai_color_scheme': 'blue',
			'ai_structure_features': 'test structure',
			'ai_template_features': 'test features',
			'ppthub_category': 'general',
			'language': '中文',
		}
		import json
		ai_output_path.write_text(json.dumps(fake_ai_output, ensure_ascii=False), encoding='utf-8')
		
		result = run_ai_enrich(conn, source_batch_id='b1', etl_out=etl_out, workshop=wc)
		assert result['ppthub_category'] == 'report'
		stage = get_stage_status(conn, aid='a1', stage=StageName.C)
		assert stage is not None and stage.status == StageStatus.success
	finally:
		conn.close()


def test_ai_fallback_and_persist(tmp_path: Path, monkeypatch) -> None:
	project_root = _setup_project_root(
		tmp_path,
		"categories:\n  a:\n    priority: 1\n    keywords:\n      - x\n",
	)

	def fake_run(args, capture_output, text, timeout=None):  # noqa: ANN001
		output_path = Path(args[-1])
		out = {
			'ai_summary': 'This is a test summary for the presentation',
			'ai_keywords': ['k1', 'k2'],
			'ai_scenario': 'business presentation',
			'ai_color_scheme': 'blue',
			'ai_structure_features': 'sections',
			'ai_template_features': 'clean',
			'ppthub_category': 'general',
			'language': '中文',
		}
		output_path.parent.mkdir(parents=True, exist_ok=True)
		output_path.write_text(json.dumps(out, ensure_ascii=False), encoding='utf-8')

		class R:
			returncode = 0
			stdout = ''
			stderr = ''
		return R()

	import subprocess  # noqa: WPS433
	monkeypatch.setattr(subprocess, 'run', fake_run)

	config = WorkshopCConfig(
		project_root=project_root,
		prompt_output_dir=tmp_path / 'prompts',
		ai_output_dir=tmp_path / 'ai',
		ai_command=['cli', '{prompt}', '{output}'],
	)
	wc = WorkshopC(config)

	# Create test PPTX
	pptx_path = tmp_path / 'test.pptx'
	_create_test_pptx(pptx_path)

	etl_out = EtlOutput(
		aid='a2',
		channel_id='chan',
		local_pptx_path=pptx_path,
		local_cover_path=None,
		pages_count=1,
		file_size_kb=1,
		meta={'title': '未知标题', 'original_tags': [], 'channel_name': 'test'},
	)
	db = tmp_path / 'assets.db'
	init_db(db)
	conn = sqlite3.connect(db)
	conn.row_factory = sqlite3.Row
	try:
		result = run_ai_enrich(conn, source_batch_id='b1', etl_out=etl_out, workshop=wc)
		assert result['ppthub_category'] == 'general'
		assert result['ai_meta'] is not None
		stage = get_stage_status(conn, aid='a2', stage=StageName.C)
		assert stage is not None and stage.status == StageStatus.success
	finally:
		conn.close()


def test_ai_output_invalid_raises(tmp_path: Path, monkeypatch) -> None:
	project_root = _setup_project_root(
		tmp_path,
		"categories:\n  a:\n    priority: 1\n    keywords:\n      - x\n",
	)

	def fake_run(args, capture_output, text, timeout=None):  # noqa: ANN001
		output_path = Path(args[-1])
		out = {'ai_summary': 'short'}  # Missing required fields
		output_path.parent.mkdir(parents=True, exist_ok=True)
		output_path.write_text(json.dumps(out, ensure_ascii=False), encoding='utf-8')

		class R:
			returncode = 0
			stdout = ''
			stderr = ''
		return R()

	import subprocess  # noqa: WPS433
	monkeypatch.setattr(subprocess, 'run', fake_run)

	config = WorkshopCConfig(
		project_root=project_root,
		prompt_output_dir=tmp_path / 'prompts',
		ai_output_dir=tmp_path / 'ai',
		ai_command=['cli', '{prompt}', '{output}'],
	)
	wc = WorkshopC(config)

	# Create test PPTX
	pptx_path = tmp_path / 'test.pptx'
	_create_test_pptx(pptx_path)

	etl_out = EtlOutput(
		aid='a3',
		channel_id='chan',
		local_pptx_path=pptx_path,
		local_cover_path=None,
		pages_count=1,
		file_size_kb=1,
		meta={'title': '未知标题', 'original_tags': [], 'channel_name': 'test'},
	)

	with pytest.raises(ValueError):
		wc.run(etl_out)
