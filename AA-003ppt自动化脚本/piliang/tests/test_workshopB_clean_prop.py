from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from factory.workshops.workshopB import WorkshopB, WorkshopBConfig
from factory.types import EtlOutput


def _count_forbidden(prs: Presentation, token: str) -> int:
	count = 0
	for slide in prs.slides:
		for shape in slide.shapes:
			if not shape.has_text_frame:
				continue
			if token in shape.text:
				count += 1
	return count


def _build_ppt_with_tail(tmp_path: Path, *, tail_forbidden: str) -> Path:
	from pptx import Presentation as Ppt

	prs = Ppt()
	layout = prs.slide_layouts[5]
	for i in range(6):
		s = prs.slides.add_slide(layout)
		s.shapes.title.text = f'slide-{i}'
	# add forbidden to last two slides
	prs.slides[-1].shapes.title.text = tail_forbidden
	prs.slides[-2].shapes.title.text = tail_forbidden
	raw = tmp_path / 'raw.pptx'
	prs.save(raw)
	return raw


def test_property6_tail_clean_and_brand_unique(tmp_path: Path) -> None:
	token = 'BADWORD'
	raw = _build_ppt_with_tail(tmp_path, tail_forbidden=token)

	brand_path = tmp_path / 'brand.pptx'
	from pptx import Presentation as Ppt

	brand = Ppt()
	slide = brand.slides.add_slide(brand.slide_layouts[5])
	slide.shapes.title.text = 'BRAND END'
	brand.save(brand_path)

	etl_out = EtlOutput(
		aid='a1',
		channel_id='chan',
		local_pptx_path=raw,
		local_cover_path=None,
		pages_count=6,
		file_size_kb=1,
		meta={},
	)

	config = WorkshopBConfig(
		output_dir=tmp_path / 'clean',
		forbidden_keywords=[token],
		brand_end_slide_path=brand_path,
	)
	wb = WorkshopB(config)

	clean = wb.deep_clean(etl_out)
	prs = Presentation(clean.clean_pptx_path)

	# Property: no forbidden tokens in content
	assert _count_forbidden(prs, token) == 0
	# Property: brand slide present exactly once (as last slide)
	texts = [shape.text for slide in prs.slides for shape in slide.shapes if getattr(shape, 'has_text_frame', False)]
	assert texts.count('BRAND END') == 1
	assert texts[-1] == 'BRAND END'

	# Property: tail removal occurred (original slides 6 -> expect <=3 before brand)
	assert len(prs.slides) <= 4  # 3 remaining + brand
