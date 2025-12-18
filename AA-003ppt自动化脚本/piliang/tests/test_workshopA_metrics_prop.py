"""Property 4: 车间 A 基础指标与命名不变量

For any 可被 python-pptx 正常解析的 PPTX 文件，车间 A 产出：
- pages_count == len(slides)
- file_size_kb == ceil(byte_size/1024)
- local_pptx_path.name 必然包含 aid 子串

Validates: Requirements 4.1, 4.2
"""
from __future__ import annotations

import json
import math
from pathlib import Path

from pptx import Presentation

from factory.workshops.workshopA import WorkshopA, WorkshopAConfig
from factory.types import StandardInputPackage


def test_property4_pages_count_and_size(tmp_path: Path) -> None:
	# Create minimal PPTX for test
	input_dir = tmp_path / 'input'
	input_dir.mkdir(parents=True, exist_ok=True)
	pptx_path = input_dir / 'main.pptx'
	prs = Presentation()
	for _ in range(5):
		prs.slides.add_slide(prs.slide_layouts[6])
	prs.save(pptx_path)

	# Create meta.json
	meta_path = input_dir / 'meta.json'
	meta_path.write_text(json.dumps({'title': 'Test Title', 'original_tags': []}), encoding='utf-8')

	# Get expected values
	expected_pages = 5
	expected_size_kb = math.ceil(pptx_path.stat().st_size / 1024)

	# Run workshop A
	config = WorkshopAConfig(output_dir=tmp_path / 'output')
	wa = WorkshopA(config)
	pkg = StandardInputPackage(
		aid='test123',
		channel_id='test_channel',
		root_dir=input_dir,
		main_pptx_path=pptx_path,
		meta_path=meta_path,
		cover_path=None,
	)
	result = wa.etl(pkg)

	# Verify Property 4
	assert result.pages_count == expected_pages
	assert result.file_size_kb == expected_size_kb
	assert 'test123' in result.local_pptx_path.name


def test_property4_filename_contains_aid(tmp_path: Path) -> None:
	"""Verify that output filename always contains the aid."""
	# Create minimal PPTX
	input_dir = tmp_path / 'input'
	input_dir.mkdir(parents=True, exist_ok=True)
	pptx_path = input_dir / 'main.pptx'
	prs = Presentation()
	prs.slides.add_slide(prs.slide_layouts[6])
	prs.save(pptx_path)

	test_cases = [
		('12345', '测试标题'),
		('abc', 'Test Title'),
		('139646', '灰色简约曲线背景'),
	]

	for aid, title in test_cases:
		meta_path = input_dir / 'meta.json'
		meta_path.write_text(json.dumps({'title': title}), encoding='utf-8')

		config = WorkshopAConfig(output_dir=tmp_path / 'output' / aid)
		wa = WorkshopA(config)
		pkg = StandardInputPackage(
			aid=aid,
			channel_id='test',
			root_dir=input_dir,
			main_pptx_path=pptx_path,
			meta_path=meta_path,
			cover_path=None,
		)
		result = wa.etl(pkg)
		assert aid in result.local_pptx_path.name, f'aid {aid} not in {result.local_pptx_path.name}'
