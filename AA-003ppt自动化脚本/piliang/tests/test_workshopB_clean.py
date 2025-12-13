from __future__ import annotations

import sqlite3
from pathlib import Path

from pptx import Presentation

from factory.db.init_db import init_db
from factory.db.dao import get_stage_status
from factory.stages import StageName, StageStatus
from factory.types import EtlOutput
from factory.workshops.workshopB import WorkshopB, WorkshopBConfig, run_clean


def _make_fixture_pptx(path: Path, *, text: str, title: str) -> None:
	from pptx import Presentation as Ppt

	prs = Ppt()
	slide_layout = prs.slide_layouts[5]  # title only layout
	slide = prs.slides.add_slide(slide_layout)
	title_shape = slide.shapes.title
	title_shape.text = text
	prs.core_properties.title = title
	prs.save(path)


def test_workshopB_removes_forbidden_and_writes_stage(tmp_path: Path) -> None:
	pptx_path = tmp_path / 'raw.pptx'
	_make_fixture_pptx(pptx_path, text='hello 第一PPT world', title='第一PPT 标题')

	etl_out = EtlOutput(
		aid='a1',
		channel_id='chan',
		local_pptx_path=pptx_path,
		local_cover_path=None,
		pages_count=1,
		file_size_kb=1,
		meta={'title': '第一PPT 标题'},
	)

	config = WorkshopBConfig(output_dir=tmp_path / 'clean', forbidden_keywords=['第一PPT'])
	wb = WorkshopB(config)

	db_path = tmp_path / 'assets.db'
	init_db(db_path)
	conn = sqlite3.connect(db_path)
	conn.row_factory = sqlite3.Row
	try:
		clean_out = run_clean(conn, source_batch_id='b1', etl_out=etl_out, workshop=wb)
		assert clean_out.clean_pptx_path.exists()
		assert clean_out.forbidden_keywords == ['第一PPT']
		assert 'WARN_FORBIDDEN_REPLACED' in clean_out.warnings

		# Verify PPTX content no longer contains forbidden token.
		prs = Presentation(clean_out.clean_pptx_path)
		texts = []
		for slide in prs.slides:
			for shape in slide.shapes:
				if shape.has_text_frame:
					texts.append(shape.text)
		assert all('第一PPT' not in t for t in texts)
		assert '第一PPT' not in prs.core_properties.title

		stage = get_stage_status(conn, aid='a1', stage=StageName.B)
		assert stage is not None
		assert stage.status == StageStatus.success
	finally:
		conn.close()


def test_workshopB_tail_prune_and_brand_insert(tmp_path: Path) -> None:
	# Build pptx with forbidden tokens in tail slides.
	from pptx import Presentation as Ppt

	prs = Ppt()
	layout = prs.slide_layouts[5]
	for i in range(5):
		s = prs.slides.add_slide(layout)
		s.shapes.title.text = f'slide-{i}'
	s_last = prs.slides[-1]
	s_last.shapes.title.text = 'BAD 第一PPT'
	prs.core_properties.title = 'Header'
	raw_path = tmp_path / 'raw2.pptx'
	prs.save(raw_path)

	# brand end slide fixture
	brand_path = tmp_path / 'brand.pptx'
	_make_fixture_pptx(brand_path, text='BRAND END', title='Brand')

	etl_out = EtlOutput(
		aid='a2',
		channel_id='chan',
		local_pptx_path=raw_path,
		local_cover_path=None,
		pages_count=5,
		file_size_kb=1,
		meta={'title': 'Header'},
	)

	config = WorkshopBConfig(
		output_dir=tmp_path / 'clean',
		forbidden_keywords=['第一PPT'],
		brand_end_slide_path=brand_path,
	)
	wb = WorkshopB(config)

	db_path = tmp_path / 'assets.db'
	init_db(db_path)
	conn = sqlite3.connect(db_path)
	conn.row_factory = sqlite3.Row
	try:
		clean_out = run_clean(conn, source_batch_id='b1', etl_out=etl_out, workshop=wb)
		prs_clean = Presentation(clean_out.clean_pptx_path)
		# removed last 3 + forbidden tail -> expect <=3 slides; brand appended as last.
		texts = [shape.text for slide in prs_clean.slides for shape in slide.shapes if getattr(shape, 'has_text_frame', False)]
		assert all('第一PPT' not in t for t in texts)
		assert 'BRAND END' in texts[-1]
		assert 'WARN_TAIL_REMOVED' in clean_out.warnings
		assert 'WARN_BRAND_APPENDED' in clean_out.warnings
		stage = get_stage_status(conn, aid='a2', stage=StageName.B)
		assert stage is not None and stage.status == StageStatus.success
	finally:
		conn.close()


def test_workshopB_fails_on_invalid_output(tmp_path: Path) -> None:
	# Force invalid pptx after cleaning to ensure validator catches it.
	raw = tmp_path / 'raw3.pptx'
	_make_fixture_pptx(raw, text='ok', title='ok')

	class BrokenWB(WorkshopB):
		def deep_clean(self, etl_out):  # type: ignore[override]
			co = super().deep_clean(etl_out)
			co.clean_pptx_path.write_bytes(b'invalid pptx content')
			return co

	config = WorkshopBConfig(output_dir=tmp_path / 'clean', forbidden_keywords=['X'])
	wb = BrokenWB(config)
	etl_out = EtlOutput(
		aid='a3',
		channel_id='chan',
		local_pptx_path=raw,
		local_cover_path=None,
		pages_count=1,
		file_size_kb=1,
		meta={},
	)

	db_path = tmp_path / 'assets.db'
	init_db(db_path)
	conn = sqlite3.connect(db_path)
	conn.row_factory = sqlite3.Row
	try:
		try:
			run_clean(conn, source_batch_id='b1', etl_out=etl_out, workshop=wb)
			assert False, 'expected failure'
		except Exception:
			pass
		stage = get_stage_status(conn, aid='a3', stage=StageName.B)
		assert stage is not None
		assert stage.status == StageStatus.failed
	finally:
		conn.close()
