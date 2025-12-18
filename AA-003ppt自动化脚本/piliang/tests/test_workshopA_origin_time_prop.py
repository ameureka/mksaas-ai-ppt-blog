"""Property 5: 来源时间口径稳定传递

For any 包含 origin_updated_at 的资产，车间 A/D/F 在导出 created_at 口径时
保持同一来源时间（或在缺失时保持旧值），重跑不改变已写入的来源时间。

Validates: Requirements 4.3
"""
from __future__ import annotations

import json
from datetime import datetime
from pathlib import Path

from pptx import Presentation

from factory.workshops.workshopA import WorkshopA, WorkshopAConfig
from factory.types import StandardInputPackage


def test_property5_origin_time_preserved(tmp_path: Path) -> None:
	"""Verify origin_updated_at is correctly parsed and preserved."""
	input_dir = tmp_path / 'input'
	input_dir.mkdir(parents=True, exist_ok=True)
	pptx_path = input_dir / 'main.pptx'
	prs = Presentation()
	prs.slides.add_slide(prs.slide_layouts[6])
	prs.save(pptx_path)

	# Test with ISO format
	meta_path = input_dir / 'meta.json'
	meta_path.write_text(json.dumps({
		'title': 'Test',
		'origin_updated_at': '2025-10-01T12:30:00',
	}), encoding='utf-8')

	config = WorkshopAConfig(output_dir=tmp_path / 'output')
	wa = WorkshopA(config)
	pkg = StandardInputPackage(
		aid='test1',
		channel_id='test',
		root_dir=input_dir,
		main_pptx_path=pptx_path,
		meta_path=meta_path,
		cover_path=None,
	)
	result = wa.etl(pkg)

	assert result.origin_updated_at is not None
	assert result.origin_updated_at.year == 2025
	assert result.origin_updated_at.month == 10
	assert result.origin_updated_at.day == 1


def test_property5_origin_time_missing_returns_none(tmp_path: Path) -> None:
	"""Verify missing origin_updated_at returns None."""
	input_dir = tmp_path / 'input'
	input_dir.mkdir(parents=True, exist_ok=True)
	pptx_path = input_dir / 'main.pptx'
	prs = Presentation()
	prs.slides.add_slide(prs.slide_layouts[6])
	prs.save(pptx_path)

	meta_path = input_dir / 'meta.json'
	meta_path.write_text(json.dumps({'title': 'Test'}), encoding='utf-8')

	config = WorkshopAConfig(output_dir=tmp_path / 'output')
	wa = WorkshopA(config)
	pkg = StandardInputPackage(
		aid='test2',
		channel_id='test',
		root_dir=input_dir,
		main_pptx_path=pptx_path,
		meta_path=meta_path,
		cover_path=None,
	)
	result = wa.etl(pkg)

	assert result.origin_updated_at is None


def test_property5_rerun_preserves_time(tmp_path: Path) -> None:
	"""Verify rerun produces same origin_updated_at."""
	input_dir = tmp_path / 'input'
	input_dir.mkdir(parents=True, exist_ok=True)
	pptx_path = input_dir / 'main.pptx'
	prs = Presentation()
	prs.slides.add_slide(prs.slide_layouts[6])
	prs.save(pptx_path)

	meta_path = input_dir / 'meta.json'
	meta_path.write_text(json.dumps({
		'title': 'Test',
		'origin_updated_at': '2025-06-15T08:00:00',
	}), encoding='utf-8')

	config = WorkshopAConfig(output_dir=tmp_path / 'output')
	wa = WorkshopA(config)
	pkg = StandardInputPackage(
		aid='test3',
		channel_id='test',
		root_dir=input_dir,
		main_pptx_path=pptx_path,
		meta_path=meta_path,
		cover_path=None,
	)

	# Run twice
	result1 = wa.etl(pkg)
	result2 = wa.etl(pkg)

	assert result1.origin_updated_at == result2.origin_updated_at
