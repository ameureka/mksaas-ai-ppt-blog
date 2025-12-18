"""
AI Enrichment 端到端集成测试
完整流程测试（mock AI）
"""

from __future__ import annotations

import json
import sqlite3
from pathlib import Path
from typing import Any
from unittest.mock import MagicMock, patch

import pytest
from pptx import Presentation
from pptx.util import Inches

from src.factory.ai import (
	AIEnrichmentService,
	AiEnrichmentOutput,
	BatchProcessReport,
	EnrichmentResult,
	EnrichmentServiceConfig,
	EtlOutput,
	create_etl_output,
)
from src.factory.ai.config_loader import AIEnrichmentConfig


@pytest.fixture
def project_root(tmp_path: Path) -> Path:
	"""创建项目根目录结构"""
	# 创建配置目录
	configs_dir = tmp_path / 'configs'
	configs_dir.mkdir()

	# 创建 forbidden-keywords.txt
	forbidden_file = configs_dir / 'forbidden-keywords.txt'
	forbidden_file.write_text('第一PPT\n1ppt\nwww.1ppt.com\nBAD_WORD\n', encoding='utf-8')

	# 创建模板目录
	templates_dir = tmp_path / 'templates'
	templates_dir.mkdir()

	# 创建 AI Prompt 模板
	template_file = templates_dir / 'ai_prompt_template.md'
	template_file.write_text(
		'''# AI PPT 内容分析任务

## 输入信息
- 标题: {title}
- 原始标签: {original_tags}
- 频道: {channel_name}

## PPT 文本内容
{extracted_text}

## 输出要求
请输出 JSON 格式，包含以下字段：
- ai_summary: 简短摘要（50-200字）
- ai_keywords: 关键词列表（3-8个）
- ai_scenario: 使用场景
- ai_color_scheme: 配色方案
- ai_structure_features: 结构特点
- ai_template_features: 模板特点
- ppthub_category: 分类（business/education/general）
- language: 语言（中文/English/其他）
''',
		encoding='utf-8',
	)

	# 创建 AI 任务目录
	(tmp_path / 'ai_tasks' / 'pending').mkdir(parents=True)
	(tmp_path / 'ai_tasks' / 'completed').mkdir(parents=True)

	return tmp_path


@pytest.fixture
def ai_config(project_root: Path) -> AIEnrichmentConfig:
	"""创建 AI 配置"""
	return AIEnrichmentConfig(
		project_root=project_root,
		forbidden_keywords=['第一PPT', '1ppt', 'www.1ppt.com', 'BAD_WORD'],
		valid_categories={'business', 'education', 'general', 'technology'},
		category_mapping=None,
		fallback_category='general',
		ai_command_template=['echo', 'mock'],
		ai_timeout_seconds=5,
		ai_max_retries=1,
		ai_concurrency=1,
		prompt_template_path=project_root / 'templates' / 'ai_prompt_template.md',
		ai_pending_dir=project_root / 'ai_tasks' / 'pending',
		ai_completed_dir=project_root / 'ai_tasks' / 'completed',
		warnings=[],
	)


@pytest.fixture
def sample_pptx(tmp_path: Path) -> Path:
	"""创建示例 PPTX 文件"""
	pptx_path = tmp_path / 'sample.pptx'
	prs = Presentation()

	# 添加幻灯片
	slide_layout = prs.slide_layouts[6]  # 空白布局

	# 第一页
	slide1 = prs.slides.add_slide(slide_layout)
	txBox1 = slide1.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
	txBox1.text_frame.text = '商务年终总结报告'

	# 第二页
	slide2 = prs.slides.add_slide(slide_layout)
	txBox2 = slide2.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
	txBox2.text_frame.text = '本年度工作回顾\n主要成果展示\n未来规划'

	# 第三页
	slide3 = prs.slides.add_slide(slide_layout)
	txBox3 = slide3.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
	txBox3.text_frame.text = 'Thank you for your attention'

	prs.save(pptx_path)
	return pptx_path


@pytest.fixture
def sample_etl_output(sample_pptx: Path) -> EtlOutput:
	"""创建示例 EtlOutput"""
	return create_etl_output(
		aid='test-001',
		channel_id='ppt_moban',
		local_pptx_path=sample_pptx,
		meta={
			'title': '商务年终总结PPT模板',
			'tags': ['商务', '年终总结', '工作汇报'],
		},
	)


@pytest.fixture
def mock_ai_response() -> dict[str, Any]:
	"""模拟 AI 响应"""
	return {
		'ai_summary': '这是一份专业的商务年终总结PPT模板，适合企业年度工作汇报使用。',
		'ai_keywords': ['商务', '年终总结', '工作汇报', '企业', '专业'],
		'ai_scenario': '企业年度总结会议、部门工作汇报',
		'ai_color_scheme': '蓝色商务风格',
		'ai_structure_features': '包含封面、目录、内容、总结四个部分',
		'ai_template_features': '简洁大气、图文并茂',
		'ppthub_category': 'business',
		'language': '中文',
	}


class TestAIEnrichmentE2E:
	"""AI Enrichment 端到端测试"""

	@patch('src.factory.ai.enrichment_service.TextExtractor')
	def test_e2e_single_asset_skip_ai(
		self,
		mock_text_extractor_class: MagicMock,
		ai_config: AIEnrichmentConfig,
		sample_etl_output: EtlOutput,
	) -> None:
		"""测试单个 Asset 处理（跳过 AI）"""
		# Mock TextExtractor
		from src.factory.ai.text_extractor import ExtractedText

		mock_extractor = MagicMock()
		mock_extractor.extract.return_value = ExtractedText(
			full_text='商务年终总结报告\n本年度工作回顾\n主要成果展示',
			slide_texts=['商务年终总结报告', '本年度工作回顾', '主要成果展示'],
			char_count=30,
			chinese_ratio=0.9,
			english_ratio=0.05,
		)
		mock_text_extractor_class.return_value = mock_extractor

		# 创建服务
		service = AIEnrichmentService(ai_config=ai_config)

		# 处理单个 Asset（跳过 AI）
		result = service.enrich_single(sample_etl_output, skip_ai=True)

		# 验证结果
		assert result.success is True
		assert result.aid == 'test-001'
		assert result.output is not None
		assert result.output.language == '中文'
		assert result.output.ppthub_category == 'general'  # fallback
		assert len(result.output.tags_final) >= 3

	@patch('src.factory.ai.enrichment_service.TextExtractor')
	def test_e2e_single_asset_with_rule_match(
		self,
		mock_text_extractor_class: MagicMock,
		ai_config: AIEnrichmentConfig,
		sample_etl_output: EtlOutput,
	) -> None:
		"""测试单个 Asset 处理（规则命中）"""
		from src.factory.ai.text_extractor import ExtractedText

		mock_extractor = MagicMock()
		mock_extractor.extract.return_value = ExtractedText(
			full_text='商务年终总结报告',
			slide_texts=['商务年终总结报告'],
			char_count=10,
			chinese_ratio=0.9,
			english_ratio=0.05,
		)
		mock_text_extractor_class.return_value = mock_extractor

		# Mock RuleEngine
		mock_rule_engine = MagicMock()
		mock_rule_engine.match.return_value = 'business'

		# 创建服务
		service = AIEnrichmentService(ai_config=ai_config)

		# 处理单个 Asset（规则命中）
		result = service.enrich_single(
			sample_etl_output,
			rule_engine=mock_rule_engine,
			skip_ai=True,
		)

		# 验证结果
		assert result.success is True
		assert result.rule_matched is True
		assert result.output is not None
		assert result.output.ppthub_category == 'business'
		assert result.output.category_source == 'rule'

	@patch('src.factory.ai.enrichment_service.TextExtractor')
	def test_e2e_batch_processing(
		self,
		mock_text_extractor_class: MagicMock,
		ai_config: AIEnrichmentConfig,
		tmp_path: Path,
	) -> None:
		"""测试批量处理"""
		from src.factory.ai.text_extractor import ExtractedText

		mock_extractor = MagicMock()
		mock_extractor.extract.return_value = ExtractedText(
			full_text='测试内容',
			slide_texts=['测试内容'],
			char_count=10,
			chinese_ratio=0.8,
			english_ratio=0.1,
		)
		mock_text_extractor_class.return_value = mock_extractor

		# 创建多个 EtlOutput
		etl_outputs = []
		for i in range(5):
			pptx_path = tmp_path / f'test_{i}.pptx'
			pptx_path.write_bytes(b'')
			etl_outputs.append(
				create_etl_output(
					aid=f'batch-{i:03d}',
					channel_id='ppt_moban',
					local_pptx_path=pptx_path,
					meta={'title': f'测试模板{i}', 'tags': ['测试']},
				)
			)

		# 创建服务
		service = AIEnrichmentService(ai_config=ai_config)

		# 批量处理
		report = service.enrich_batch(
			batch_id='e2e-batch',
			etl_outputs=etl_outputs,
			skip_ai=True,
		)

		# 验证结果
		assert report.total == 5
		assert report.success == 5
		assert report.failed == 0
		assert report.success_rate == 1.0

	@patch('src.factory.ai.enrichment_service.TextExtractor')
	def test_e2e_batch_with_failures(
		self,
		mock_text_extractor_class: MagicMock,
		ai_config: AIEnrichmentConfig,
		tmp_path: Path,
	) -> None:
		"""测试批量处理包含失败（错误隔离）"""
		from src.factory.ai.text_extractor import ExtractedText

		call_count = {'count': 0}

		def extract_side_effect(path: Path) -> ExtractedText:
			call_count['count'] += 1
			if call_count['count'] == 3:
				raise RuntimeError('Simulated extraction error')
			return ExtractedText(
				full_text='测试内容',
				slide_texts=['测试内容'],
				char_count=10,
				chinese_ratio=0.8,
				english_ratio=0.1,
			)

		mock_extractor = MagicMock()
		mock_extractor.extract.side_effect = extract_side_effect
		mock_text_extractor_class.return_value = mock_extractor

		# 创建多个 EtlOutput
		etl_outputs = []
		for i in range(5):
			pptx_path = tmp_path / f'test_{i}.pptx'
			pptx_path.write_bytes(b'')
			etl_outputs.append(
				create_etl_output(
					aid=f'batch-{i:03d}',
					channel_id='ppt_moban',
					local_pptx_path=pptx_path,
					meta={'title': f'测试模板{i}', 'tags': ['测试']},
				)
			)

		# 创建服务
		service = AIEnrichmentService(ai_config=ai_config)

		# 批量处理
		report = service.enrich_batch(
			batch_id='e2e-batch-fail',
			etl_outputs=etl_outputs,
			skip_ai=True,
		)

		# 验证结果 - 单个失败不阻塞其他
		assert report.total == 5
		assert report.success == 4
		assert report.failed == 1
		assert 'batch-002' in report.failed_aids

	@patch('src.factory.ai.enrichment_service.TextExtractor')
	def test_e2e_forbidden_keyword_filtering(
		self,
		mock_text_extractor_class: MagicMock,
		ai_config: AIEnrichmentConfig,
		sample_etl_output: EtlOutput,
	) -> None:
		"""测试敏感词过滤"""
		from src.factory.ai.text_extractor import ExtractedText

		mock_extractor = MagicMock()
		mock_extractor.extract.return_value = ExtractedText(
			full_text='商务报告 第一PPT 下载',
			slide_texts=['商务报告 第一PPT 下载'],
			char_count=15,
			chinese_ratio=0.9,
			english_ratio=0.05,
		)
		mock_text_extractor_class.return_value = mock_extractor

		# 修改 meta 包含敏感词
		etl_with_forbidden = create_etl_output(
			aid='test-forbidden',
			channel_id='ppt_moban',
			local_pptx_path=sample_etl_output.local_pptx_path,
			meta={
				'title': '商务PPT模板',
				'tags': ['商务', '第一PPT', '年终总结', 'BAD_WORD'],
			},
		)

		# 创建服务
		service = AIEnrichmentService(ai_config=ai_config)

		# 处理
		result = service.enrich_single(etl_with_forbidden, skip_ai=True)

		# 验证敏感词被过滤
		assert result.success is True
		assert result.output is not None
		# tags_final 不应包含敏感词
		for tag in result.output.tags_final:
			assert '第一PPT' not in tag
			assert 'BAD_WORD' not in tag

	@patch('src.factory.ai.enrichment_service.TextExtractor')
	def test_e2e_language_detection(
		self,
		mock_text_extractor_class: MagicMock,
		ai_config: AIEnrichmentConfig,
		tmp_path: Path,
	) -> None:
		"""测试语言检测"""
		from src.factory.ai.text_extractor import ExtractedText

		# 测试中文内容 - 使用纯中文文本确保中文比例 > 30%
		mock_extractor = MagicMock()
		chinese_text = '商务年终总结报告模板工作汇报企业管理'  # 纯中文，比例 100%
		mock_extractor.extract.return_value = ExtractedText(
			full_text=chinese_text,
			slide_texts=[chinese_text],
			char_count=len(chinese_text),
			chinese_ratio=1.0,
			english_ratio=0.0,
		)
		mock_text_extractor_class.return_value = mock_extractor

		# 创建服务（在 mock 设置之后）
		service = AIEnrichmentService(ai_config=ai_config)

		pptx_path = tmp_path / 'chinese.pptx'
		pptx_path.write_bytes(b'')
		etl_chinese = create_etl_output(
			aid='test-chinese',
			channel_id='ppt_moban',
			local_pptx_path=pptx_path,
			meta={'title': '中文模板', 'tags': ['中文']},
		)

		result_zh = service.enrich_single(etl_chinese, skip_ai=True)
		assert result_zh.output is not None
		assert result_zh.output.language == '中文'

		# 测试英文内容 - 使用纯英文文本确保英文比例 > 70%
		english_text = 'BusinessReportTemplateForAnnualSummaryAndWorkPresentation'  # 纯英文字母
		mock_extractor.extract.return_value = ExtractedText(
			full_text=english_text,
			slide_texts=[english_text],
			char_count=len(english_text),
			chinese_ratio=0.0,
			english_ratio=1.0,
		)

		pptx_path_en = tmp_path / 'english.pptx'
		pptx_path_en.write_bytes(b'')
		etl_english = create_etl_output(
			aid='test-english',
			channel_id='ppt_moban',
			local_pptx_path=pptx_path_en,
			meta={'title': 'English Template', 'tags': ['english']},
		)

		result_en = service.enrich_single(etl_english, skip_ai=True)
		assert result_en.output is not None
		assert result_en.output.language == 'English'

	@patch('src.factory.ai.enrichment_service.TextExtractor')
	def test_e2e_tags_normalization(
		self,
		mock_text_extractor_class: MagicMock,
		ai_config: AIEnrichmentConfig,
		tmp_path: Path,
	) -> None:
		"""测试标签标准化"""
		from src.factory.ai.text_extractor import ExtractedText

		mock_extractor = MagicMock()
		mock_extractor.extract.return_value = ExtractedText(
			full_text='商务报告',
			slide_texts=['商务报告'],
			char_count=4,
			chinese_ratio=1.0,
			english_ratio=0.0,
		)
		mock_text_extractor_class.return_value = mock_extractor

		pptx_path = tmp_path / 'tags_test.pptx'
		pptx_path.write_bytes(b'')

		# 测试标签过多的情况
		etl_many_tags = create_etl_output(
			aid='test-many-tags',
			channel_id='ppt_moban',
			local_pptx_path=pptx_path,
			meta={
				'title': '商务PPT模板',
				'tags': ['tag1', 'tag2', 'tag3', 'tag4', 'tag5', 'tag6', 'tag7', 'tag8', 'tag9', 'tag10'],
			},
		)

		service = AIEnrichmentService(ai_config=ai_config)
		result = service.enrich_single(etl_many_tags, skip_ai=True)

		assert result.success is True
		assert result.output is not None
		# 标签应该被收敛到 3-8 个
		assert 3 <= len(result.output.tags_final) <= 8

	@patch('src.factory.ai.enrichment_service.TextExtractor')
	def test_e2e_description_building(
		self,
		mock_text_extractor_class: MagicMock,
		ai_config: AIEnrichmentConfig,
		sample_etl_output: EtlOutput,
	) -> None:
		"""测试描述构建"""
		from src.factory.ai.text_extractor import ExtractedText

		mock_extractor = MagicMock()
		mock_extractor.extract.return_value = ExtractedText(
			full_text='商务年终总结报告模板',
			slide_texts=['商务年终总结报告模板'],
			char_count=10,
			chinese_ratio=1.0,
			english_ratio=0.0,
		)
		mock_text_extractor_class.return_value = mock_extractor

		service = AIEnrichmentService(ai_config=ai_config)
		result = service.enrich_single(sample_etl_output, skip_ai=True)

		assert result.success is True
		assert result.output is not None
		# 描述应该被构建
		assert len(result.output.description_final) > 0
		# 描述不应超过 500 字
		assert len(result.output.description_final) <= 500


class TestAIEnrichmentServiceConfig:
	"""AIEnrichmentService 配置测试"""

	def test_default_config(self) -> None:
		"""测试默认配置"""
		config = EnrichmentServiceConfig()

		assert config.ai_timeout_seconds == 60
		assert config.ai_max_retries == 2
		assert config.max_concurrency == 1
		assert config.description_min_length == 50
		assert config.description_max_length == 500
		assert config.tags_min_count == 3
		assert config.tags_max_count == 8

	def test_custom_config(self) -> None:
		"""测试自定义配置"""
		config = EnrichmentServiceConfig(
			ai_timeout_seconds=120,
			max_concurrency=4,
			tags_max_count=10,
		)

		assert config.ai_timeout_seconds == 120
		assert config.max_concurrency == 4
		assert config.tags_max_count == 10


class TestEtlOutputCreation:
	"""EtlOutput 创建测试"""

	def test_create_etl_output_basic(self, tmp_path: Path) -> None:
		"""测试基本创建"""
		pptx_path = tmp_path / 'test.pptx'
		pptx_path.write_bytes(b'')

		etl = create_etl_output(
			aid='test-001',
			channel_id='ppt_moban',
			local_pptx_path=pptx_path,
			meta={'title': '测试', 'tags': ['tag1']},
		)

		assert etl.aid == 'test-001'
		assert etl.channel_id == 'ppt_moban'
		assert etl.local_pptx_path == pptx_path
		assert etl.meta == {'title': '测试', 'tags': ['tag1']}

	def test_create_etl_output_string_path(self, tmp_path: Path) -> None:
		"""测试字符串路径"""
		pptx_path = tmp_path / 'test.pptx'
		pptx_path.write_bytes(b'')

		etl = create_etl_output(
			aid='test-001',
			channel_id='ppt_moban',
			local_pptx_path=str(pptx_path),
		)

		assert isinstance(etl.local_pptx_path, Path)

	def test_create_etl_output_default_meta(self, tmp_path: Path) -> None:
		"""测试默认 meta"""
		pptx_path = tmp_path / 'test.pptx'
		pptx_path.write_bytes(b'')

		etl = create_etl_output(
			aid='test-001',
			channel_id='ppt_moban',
			local_pptx_path=pptx_path,
		)

		assert etl.meta == {}
