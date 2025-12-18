"""Property 14: 车间 D 成品打包目录结构与路径回填

For any 已完成车间 B/C 的 Asset：
- 车间 D 将清洗后的 PPTX 归档到 output/{channel_id}/ 目录
- 本地成品文件命名包含 aid
- 完成打包后 pptx_path 非空

Validates: Requirements 7.1, 7.2, 7.3
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from factory.workshops.workshopD import WorkshopD, WorkshopDConfig
from factory.types import CleanOutput, EtlOutput


def _create_test_pptx(path: Path) -> None:
	prs = Presentation()
	prs.slides.add_slide(prs.slide_layouts[6])
	prs.save(path)


def test_property14_output_directory_structure(tmp_path: Path) -> None:
	"""Verify output is placed in output/{channel_id}/ directory."""
	clean_dir = tmp_path / 'clean'
	clean_dir.mkdir()
	pptx_path = clean_dir / 'test.pptx'
	_create_test_pptx(pptx_path)

	config = WorkshopDConfig(output_dir=tmp_path / 'output')
	wd = WorkshopD(config)

	etl_out = EtlOutput(
		aid='12345',
		channel_id='business',
		local_pptx_path=pptx_path,
		local_cover_path=None,
		pages_count=1,
		file_size_kb=10,
		meta={},
	)
	clean_out = CleanOutput(
		aid='12345',
		channel_id='business',
		clean_pptx_path=pptx_path,
	)

	result = wd.pack(etl_out, clean_out)

	# Verify directory structure
	assert result.pptx_path.parent.name == 'business'
	assert 'output' in str(result.pptx_path.parent.parent)

def test_property14_filename_contains_aid(tmp_path: Path) -> None:
	"""Verify filename contains aid."""
	clean_dir = tmp_path / 'clean'
	clean_dir.mkdir()
	pptx_path = clean_dir / 'test.pptx'
	_create_test_pptx(pptx_path)

	config = WorkshopDConfig(output_dir=tmp_path / 'output')
	wd = WorkshopD(config)

	test_aids = ['12345', 'abc', '139646']

	for aid in test_aids:
		etl_out = EtlOutput(
			aid=aid,
			channel_id='test',
			local_pptx_path=pptx_path,
			local_cover_path=None,
			pages_count=1,
			file_size_kb=10,
			meta={},
		)
		clean_out = CleanOutput(
			aid=aid,
			channel_id='test',
			clean_pptx_path=pptx_path,
		)
		result = wd.pack(etl_out, clean_out)
		assert aid in result.pptx_path.name, f'aid {aid} not in {result.pptx_path.name}'

def test_property14_pptx_path_exists(tmp_path: Path) -> None:
	"""Verify pptx_path exists after pack."""
	clean_dir = tmp_path / 'clean'
	clean_dir.mkdir()
	pptx_path = clean_dir / 'test.pptx'
	_create_test_pptx(pptx_path)

	config = WorkshopDConfig(output_dir=tmp_path / 'output')
	wd = WorkshopD(config)

	etl_out = EtlOutput(
		aid='test',
		channel_id='test',
		local_pptx_path=pptx_path,
		local_cover_path=None,
		pages_count=1,
		file_size_kb=10,
		meta={},
	)
	clean_out = CleanOutput(
		aid='test',
		channel_id='test',
		clean_pptx_path=pptx_path,
	)

	result = wd.pack(etl_out, clean_out)

	assert result.pptx_path is not None
	assert result.pptx_path.exists()
