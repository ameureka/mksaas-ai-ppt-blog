from __future__ import annotations

import sqlite3
from pathlib import Path

from pptx import Presentation

from factory.db.init_db import init_db
from factory.db.dao import get_stage_status
from factory.stages import StageName, StageStatus
from factory.types import CleanOutput, EtlOutput
from factory.workshops.workshopD import WorkshopD, WorkshopDConfig, run_pack


def _create_test_pptx(path: Path) -> None:
	prs = Presentation()
	prs.slides.add_slide(prs.slide_layouts[6])
	prs.save(path)


def test_workshopD_packs_outputs_and_writes_stage(tmp_path: Path) -> None:
	# Create test PPTX
	src_dir = tmp_path / 'src'
	src_dir.mkdir()
	pptx_path = src_dir / 'main.pptx'
	_create_test_pptx(pptx_path)

	etl_out = EtlOutput(
		aid='a1',
		channel_id='chan',
		local_pptx_path=pptx_path,
		local_cover_path=None,
		pages_count=1,
		file_size_kb=1,
		meta={'title': 'T'},
	)
	clean_out = CleanOutput(
		aid='a1',
		channel_id='chan',
		clean_pptx_path=pptx_path,
	)

	config = WorkshopDConfig(output_dir=tmp_path / 'out')
	wd = WorkshopD(config)

	db = tmp_path / 'assets.db'
	init_db(db)
	conn = sqlite3.connect(db)
	conn.row_factory = sqlite3.Row
	try:
		result = run_pack(conn, source_batch_id='b1', etl_out=etl_out, clean_out=clean_out, workshop=wd)
		assert result.pptx_path.exists()

		stage = get_stage_status(conn, aid='a1', stage=StageName.D)
		assert stage is not None
		assert stage.status == StageStatus.success
		assert 'local_pptx_path' in stage.artifacts
	finally:
		conn.close()
