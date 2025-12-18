"""
PromptBuilder 模块
负责构建 AI Prompt 文件
"""

from __future__ import annotations

import json
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from pptx import Presentation

from .prompt_template import PromptTemplate, PromptTemplateConfig
from .text_extractor import ExtractedText, TextExtractor


def _extract_text_legacy(prs: Presentation) -> str:
	"""旧版文本提取（保持向后兼容）"""
	chunks: list[str] = []
	for slide in prs.slides:
		for shape in slide.shapes:
			if not shape.has_text_frame:
				continue
			text = shape.text
			if text:
				chunks.append(text)
	return '\n'.join(chunks)


@dataclass(frozen=True)
class PromptBuilderConfig:
	"""PromptBuilder 配置"""

	output_dir: Path
	template_path: Path | None = None
	max_text_length: int = 5000
	use_template: bool = True


@dataclass(frozen=True)
class BuiltPrompt:
	"""构建的 Prompt 结果"""

	aid: str
	prompt_path: Path
	prompt_content: str
	extracted_text: ExtractedText | None
	template_used: bool


class PromptBuilder:
	"""Prompt 构建器"""

	def __init__(self, config: PromptBuilderConfig) -> None:
		"""
		初始化 Prompt 构建器

		Args:
			config: 构建器配置
		"""
		self._config = config
		self._text_extractor = TextExtractor()
		self._prompt_template: PromptTemplate | None = None

		# 加载模板（如果配置了）
		if config.template_path and config.template_path.exists() and config.use_template:
			self._prompt_template = PromptTemplate(
				PromptTemplateConfig(template_path=config.template_path)
			)
			self._prompt_template.load()

	def build(
		self, *, aid: str, title: str, meta: dict[str, Any], pptx_path: Path
	) -> Path:
		"""
		构建 Prompt 文件（旧接口，保持向后兼容）

		Args:
			aid: Asset ID
			title: 标题
			meta: 元数据
			pptx_path: PPTX 文件路径

		Returns:
			Prompt 文件路径
		"""
		result = self.build_prompt(
			aid=aid,
			title=title,
			meta=meta,
			pptx_path=pptx_path,
			channel_name=meta.get('channel_id', 'unknown'),
		)
		return result.prompt_path

	def build_prompt(
		self,
		*,
		aid: str,
		title: str,
		meta: dict[str, Any],
		pptx_path: Path,
		channel_name: str = 'unknown',
	) -> BuiltPrompt:
		"""
		构建 Prompt 文件（增强版）

		Args:
			aid: Asset ID
			title: 标题
			meta: 元数据
			pptx_path: PPTX 文件路径
			channel_name: 频道名称

		Returns:
			BuiltPrompt 对象
		"""
		# 提取文本
		extracted_text = self._text_extractor.extract(pptx_path)

		# 获取原始标签
		original_tags = meta.get('tags', [])
		if isinstance(original_tags, str):
			original_tags = [original_tags]

		# 截断文本
		text_content = extracted_text.full_text
		if len(text_content) > self._config.max_text_length:
			text_content = text_content[: self._config.max_text_length]

		# 构建 Prompt 内容
		template_used = False
		if self._prompt_template and self._config.use_template:
			# 使用模板渲染
			prompt_content = self._prompt_template.render(
				title=title,
				original_tags=original_tags,
				extracted_text=text_content,
				channel_name=channel_name,
			)
			template_used = True
			out_filename = 'prompt.md'
		else:
			# 使用 JSON 格式（旧版）
			payload = {
				'aid': aid,
				'title': title,
				'meta': meta,
				'extracted_text': text_content,
			}
			prompt_content = json.dumps(payload, ensure_ascii=False, indent=2)
			out_filename = 'prompt.json'

		# 写入文件
		out_dir = self._config.output_dir / aid
		out_dir.mkdir(parents=True, exist_ok=True)
		out_path = out_dir / out_filename
		out_path.write_text(prompt_content, encoding='utf-8')

		return BuiltPrompt(
			aid=aid,
			prompt_path=out_path,
			prompt_content=prompt_content,
			extracted_text=extracted_text,
			template_used=template_used,
		)

	def build_from_text(
		self,
		*,
		aid: str,
		title: str,
		original_tags: list[str],
		extracted_text: str,
		channel_name: str = 'unknown',
	) -> BuiltPrompt:
		"""
		从已提取的文本构建 Prompt（不读取 PPTX）

		Args:
			aid: Asset ID
			title: 标题
			original_tags: 原始标签
			extracted_text: 已提取的文本
			channel_name: 频道名称

		Returns:
			BuiltPrompt 对象
		"""
		# 截断文本
		text_content = extracted_text
		if len(text_content) > self._config.max_text_length:
			text_content = text_content[: self._config.max_text_length]

		# 构建 Prompt 内容
		template_used = False
		if self._prompt_template and self._config.use_template:
			prompt_content = self._prompt_template.render(
				title=title,
				original_tags=original_tags,
				extracted_text=text_content,
				channel_name=channel_name,
			)
			template_used = True
			out_filename = 'prompt.md'
		else:
			payload = {
				'aid': aid,
				'title': title,
				'tags': original_tags,
				'extracted_text': text_content,
			}
			prompt_content = json.dumps(payload, ensure_ascii=False, indent=2)
			out_filename = 'prompt.json'

		# 写入文件
		out_dir = self._config.output_dir / aid
		out_dir.mkdir(parents=True, exist_ok=True)
		out_path = out_dir / out_filename
		out_path.write_text(prompt_content, encoding='utf-8')

		return BuiltPrompt(
			aid=aid,
			prompt_path=out_path,
			prompt_content=prompt_content,
			extracted_text=None,
			template_used=template_used,
		)
