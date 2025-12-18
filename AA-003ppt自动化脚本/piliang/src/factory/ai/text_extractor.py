"""
TextExtractor 模块
负责从 PPTX 文件中提取文本内容并进行语言分析
"""

from __future__ import annotations

import re
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation


@dataclass(frozen=True)
class ExtractedText:
	"""提取的文本数据"""

	full_text: str
	slide_texts: list[str]
	char_count: int
	chinese_ratio: float
	english_ratio: float


class TextExtractor:
	"""PPTX 文本提取器"""

	# 常见的英文占位符模式
	_PLACEHOLDER_PATTERNS = [
		r'lorem\s+ipsum',
		r'click\s+to\s+add',
		r'click\s+to\s+edit',
		r'add\s+your\s+text',
		r'type\s+here',
		r'insert\s+text',
		r'sample\s+text',
		r'placeholder',
		r'your\s+text\s+here',
		r'enter\s+text',
	]

	def __init__(self) -> None:
		self._placeholder_regex = re.compile(
			'|'.join(self._PLACEHOLDER_PATTERNS), re.IGNORECASE
		)

	def extract(self, pptx_path: Path) -> ExtractedText:
		"""
		从 PPTX 文件中提取全部可见文本

		Args:
			pptx_path: PPTX 文件路径

		Returns:
			ExtractedText 对象

		Raises:
			FileNotFoundError: PPTX 文件不存在
			ValueError: PPTX 文件无法解析
		"""
		if not pptx_path.exists():
			raise FileNotFoundError(f'PPTX file not found: {pptx_path}')

		try:
			prs = Presentation(str(pptx_path))
		except Exception as exc:
			raise ValueError(f'Failed to parse PPTX file: {pptx_path}') from exc

		slide_texts: list[str] = []
		for slide in prs.slides:
			slide_text_parts: list[str] = []
			for shape in slide.shapes:
				if not shape.has_text_frame:
					continue
				text = shape.text.strip()
				if text:
					# 过滤占位符
					filtered = self.filter_placeholders(text)
					if filtered:
						slide_text_parts.append(filtered)

			if slide_text_parts:
				slide_texts.append('\n'.join(slide_text_parts))

		full_text = '\n\n'.join(slide_texts)
		char_count = len(full_text)

		# 计算中文和英文字符占比
		chinese_count = sum(1 for c in full_text if '\u4e00' <= c <= '\u9fff')
		english_count = sum(1 for c in full_text if c.isalpha() and c.isascii())

		chinese_ratio = chinese_count / char_count if char_count > 0 else 0.0
		english_ratio = english_count / char_count if char_count > 0 else 0.0

		return ExtractedText(
			full_text=full_text,
			slide_texts=slide_texts,
			char_count=char_count,
			chinese_ratio=chinese_ratio,
			english_ratio=english_ratio,
		)

	def filter_placeholders(self, text: str) -> str:
		"""
		过滤占位符文本

		Args:
			text: 原始文本

		Returns:
			过滤后的文本
		"""
		# 移除匹配占位符模式的行
		lines = text.split('\n')
		filtered_lines: list[str] = []

		for line in lines:
			line_stripped = line.strip()
			if not line_stripped:
				continue
			# 检查是否为占位符
			if self._placeholder_regex.search(line_stripped):
				continue
			filtered_lines.append(line)

		return '\n'.join(filtered_lines)
