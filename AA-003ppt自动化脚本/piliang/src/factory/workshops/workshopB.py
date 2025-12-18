from __future__ import annotations

import io
import math
import re
import zipfile
from copy import deepcopy
from dataclasses import dataclass, field
from datetime import datetime, timezone
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from ..db.dao import record_stage, upsert_processed_asset
from ..stages import StageName, StageStatus
from ..types import CleanOutput, EtlOutput, StageRecord


@dataclass(frozen=True)
class WorkshopBConfig:
	output_dir: Path
	forbidden_keywords: list[str]
	brand_replacements: list[tuple[str, str]] = field(default_factory=list)
	brand_end_slide_path: Path | None = None  # reserved for 6.2


_SEP_PATTERN = r"[\s·．。\.\\-_/:：]*"


def _compile_literal_pattern(token: str) -> tuple[re.Pattern[str], str]:
	"""
	Compile a "fuzzy literal" token pattern:
	- matches token case-insensitively (for ASCII letters)
	- allows common separators between characters (spaces, dots, hyphens, underscores, etc.)
	- for domain-like tokens (containing dots), also matches optional http/https prefix
	"""
	parts: list[str] = []
	for ch in token:
		if ch in {'.', '。', '．', '·'}:
			parts.append(rf"{_SEP_PATTERN}[\\.。．·]{_SEP_PATTERN}")
		elif ch.isspace():
			parts.append(r"\s+")
		else:
			parts.append(re.escape(ch))

	pattern = _SEP_PATTERN.join(parts)
	if any(ch in {'.', '。', '．', '·'} for ch in token):
		pattern = rf"(?:https?://)?{pattern}"

	flags = re.IGNORECASE if any('A' <= c <= 'Z' or 'a' <= c <= 'z' for c in token) else 0
	return re.compile(pattern, flags=flags), token


def _compile_rules(
	*,
	brand_replacements: list[tuple[str, str]],
	forbidden_keywords: list[str],
) -> list[tuple[re.Pattern[str], str]]:
	"""
	Build ordered replacement rules: apply brand replacements first, then delete any remaining forbidden tokens.
	Longer tokens run first to avoid partial matches (e.g. `www.1ppt.com` before `1ppt`).
	"""
	seen_sources: set[str] = set()
	ordered_sources: list[tuple[str, str]] = []
	for source, repl in brand_replacements:
		src = (source or '').strip()
		if not src:
			continue
		if src in seen_sources:
			continue
		seen_sources.add(src)
		ordered_sources.append((src, repl))

	for kw in forbidden_keywords:
		src = (kw or '').strip()
		if not src:
			continue
		if src in seen_sources:
			continue
		seen_sources.add(src)
		ordered_sources.append((src, ''))

	# long match first (stable)
	ordered_sources.sort(key=lambda item: len(item[0]), reverse=True)

	compiled: list[tuple[re.Pattern[str], str]] = []
	for source, repl in ordered_sources:
		pat, _ = _compile_literal_pattern(source)
		compiled.append((pat, repl))
	return compiled


def _apply_rules(text: str, rules: list[tuple[re.Pattern[str], str]]) -> tuple[str, int]:
	updated = text
	total = 0
	for pattern, repl in rules:
		updated, n = pattern.subn(repl, updated)
		total += n
	return updated, total


def _iter_shapes(shapes) -> Iterable:
	"""
	Yield all shapes recursively.

	PPT templates often put watermark text inside group shapes; iterating only the
	top-level `slide.shapes` will miss those nested text frames.
	"""
	for shape in shapes:
		yield shape
		if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
			yield from _iter_shapes(shape.shapes)


def _clean_text_runs(slide, rules: list[tuple[re.Pattern[str], str]], *, hit_counter: dict[str, int]) -> None:
	for shape in _iter_shapes(slide.shapes):
		if not shape.has_text_frame:
			if getattr(shape, 'has_table', False):
				table = shape.table
				for row in table.rows:
					for cell in row.cells:
						for paragraph in cell.text_frame.paragraphs:
							original = paragraph.text or ''
							cleaned, n = _apply_rules(original, rules)
							if cleaned != original:
								hit_counter['replaced'] += max(1, n)
								paragraph.text = cleaned
			continue
		# Replace tokens at paragraph level so we also catch tokens split across runs.
		# This trades away run-level formatting for correctness in watermark removal.
		for paragraph in shape.text_frame.paragraphs:
			original = paragraph.text or ''
			cleaned, n = _apply_rules(original, rules)
			if cleaned != original:
				hit_counter['replaced'] += max(1, n)
				paragraph.text = cleaned


def _clean_core_properties(
	prs: Presentation,
	rules: list[tuple[re.Pattern[str], str]],
	*,
	hit_counter: dict[str, int],
) -> None:
	props = prs.core_properties
	fields = ['title', 'subject', 'keywords', 'comments', 'last_modified_by', 'author', 'category']
	for field in fields:
		original = getattr(props, field, None)
		if not isinstance(original, str):
			continue
		cleaned, n = _apply_rules(original, rules)
		if cleaned != original:
			hit_counter['replaced'] += max(1, n)
		setattr(props, field, cleaned)


def _clean_part_xml_blob(part, rules: list[tuple[re.Pattern[str], str]], *, hit_counter: dict[str, int]) -> None:
	"""
	Apply replacement rules directly on an XML part blob.

	Some branding appears in theme/app properties where python-pptx doesn't provide
	a high-level API. We rewrite the XML blob in-place using UTF-8.
	"""
	try:
		original = part.blob.decode('utf-8')
	except Exception:
		return
	cleaned, n = _apply_rules(original, rules)
	if cleaned == original:
		return
	hit_counter['replaced'] += max(1, n)
	part.blob = cleaned.encode('utf-8')

def _scrub_pptx_zip_xml(pptx_path: Path, rules: list[tuple[re.Pattern[str], str]], *, hit_counter: dict[str, int]) -> None:
	"""
	Scrub XML/RELS parts at the ZIP level.

	python-pptx does not expose relationship parts (`*.xml.rels`) as normal parts,
	so brand links can survive unless we rewrite those files directly.
	"""
	tmp_path = pptx_path.with_name(f'{pptx_path.name}.tmp')
	with zipfile.ZipFile(pptx_path, 'r') as zin, zipfile.ZipFile(
		tmp_path,
		'w',
		compression=zipfile.ZIP_DEFLATED,
	) as zout:
		for info in zin.infolist():
			data = zin.read(info.filename)
			if info.filename.endswith(('.xml', '.rels')):
				try:
					text = data.decode('utf-8')
				except Exception:
					text = data.decode('utf-8', errors='ignore')
				cleaned, n = _apply_rules(text, rules)
				if cleaned != text:
					hit_counter['replaced'] += max(1, n)
					data = cleaned.encode('utf-8')
			zout.writestr(info, data)
	tmp_path.replace(pptx_path)


class WorkshopB:
	def __init__(self, config: WorkshopBConfig) -> None:
		self._config = config
		self._text_rules = _compile_rules(
			brand_replacements=config.brand_replacements,
			forbidden_keywords=config.forbidden_keywords,
		)
		self._forbidden_patterns = [_compile_literal_pattern(kw)[0] for kw in config.forbidden_keywords if kw]

	def _slide_contains_forbidden(self, slide) -> bool:
		for shape in _iter_shapes(slide.shapes):
			if shape.has_text_frame:
				text = shape.text or ''
				for pat in self._forbidden_patterns:
					if pat.search(text):
						return True
			elif getattr(shape, 'has_table', False):
				table = shape.table
				for row in table.rows:
					for cell in row.cells:
						text = cell.text_frame.text or ''
						for pat in self._forbidden_patterns:
							if pat.search(text):
								return True
		return False

	@staticmethod
	def _remove_slide(prs: Presentation, index: int) -> None:
		# python-pptx does not officially support deleting slides. This approach
		# removes the slide ID element and drops the underlying relationship to
		# avoid generating duplicate part names when saving.
		sldIdLst = prs.slides._sldIdLst  # type: ignore[attr-defined]
		sldId = sldIdLst[index]
		prs.part.drop_rel(sldId.rId)  # type: ignore[attr-defined]
		sldIdLst.remove(sldId)

	@staticmethod
	def _append_brand_slide(prs: Presentation, brand_path: Path) -> None:
		if not brand_path.exists():
			return
		brand_prs = Presentation(str(brand_path))
		if not brand_prs.slides:
			return
		source_slide = brand_prs.slides[0]
		blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
		new_slide = prs.slides.add_slide(blank_layout)

		# Remove placeholder shapes in the new blank slide.
		for shape in list(new_slide.shapes):
			new_slide.shapes._spTree.remove(shape._element)  # type: ignore[attr-defined]

		# Deep copy shapes from the template slide.
		for shape in source_slide.shapes:
			new_slide.shapes._spTree.insert_element_before(deepcopy(shape._element), 'p:extLst')  # type: ignore[attr-defined]

		# Copy image relationships referenced by the copied shape XML.
		# Without this, pictures end up with dangling r:embed values and show as broken.
		embed_attr = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed'
		rid_map: dict[str, str] = {}
		for blip in new_slide._element.xpath('.//a:blip[@r:embed]'):  # type: ignore[attr-defined]
			old_rid = blip.get(embed_attr)
			if not old_rid:
				continue
			if old_rid not in rid_map:
				try:
					part = source_slide.part.related_part(old_rid)
				except KeyError:
					continue
				blob = getattr(part, 'blob', None)
				if not isinstance(blob, (bytes, bytearray)) or not blob:
					continue
				_, new_rid = new_slide.part.get_or_add_image_part(io.BytesIO(blob))
				rid_map[old_rid] = new_rid
			blip.set(embed_attr, rid_map[old_rid])

	def deep_clean(self, etl_out: EtlOutput) -> CleanOutput:
		"""
		Remove forbidden keywords from slides and core properties.
		Remove tail slides (last 3 + trailing forbidden) and append brand end slide if provided.
		"""
		if not etl_out.local_pptx_path.exists():
			raise FileNotFoundError(f'PPTX not found: {etl_out.local_pptx_path}')

		prs = Presentation(etl_out.local_pptx_path)
		hits = {'replaced': 0}

		# Collect forbidden markers per slide before replacement for tail pruning.
		slide_flags = [
			self._slide_contains_forbidden(slide) for slide in prs.slides
		]

		removed = 0
		# Remove last 3 slides.
		for _ in range(min(3, len(slide_flags))):
			self._remove_slide(prs, len(prs.slides) - 1)
			slide_flags.pop()
			removed += 1

		# Remove trailing slides still containing forbidden keywords.
		while slide_flags and slide_flags[-1]:
			self._remove_slide(prs, len(prs.slides) - 1)
			slide_flags.pop()
			removed += 1

		for slide in prs.slides:
			_clean_text_runs(slide, self._text_rules, hit_counter=hits)
		_clean_core_properties(prs, self._text_rules, hit_counter=hits)

		# Watermarks often live in masters/layouts. Clean them too so they don't
		# keep appearing across slides.
		for master in prs.slide_masters:
			_clean_text_runs(master, self._text_rules, hit_counter=hits)
		for layout in prs.slide_layouts:
			_clean_text_runs(layout, self._text_rules, hit_counter=hits)

		# Also scrub theme and app.xml blobs for stray branding (e.g. theme name).
		for part in prs.part.package.iter_parts():
			partname = str(part.partname)
			if partname.startswith('/ppt/theme/theme') and partname.endswith('.xml'):
				_clean_part_xml_blob(part, self._text_rules, hit_counter=hits)
			elif partname.endswith('.rels'):
				# Watermarks/brand links can hide in relationship targets (e.g. slideLayout rels).
				_clean_part_xml_blob(part, self._text_rules, hit_counter=hits)
			elif partname == '/docProps/app.xml':
				_clean_part_xml_blob(part, self._text_rules, hit_counter=hits)

		output_dir = self._config.output_dir / etl_out.channel_id
		output_dir.mkdir(parents=True, exist_ok=True)
		output_path = output_dir / f'{etl_out.aid}-clean.pptx'

		brand_inserted = False
		if self._config.brand_end_slide_path:
			self._append_brand_slide(prs, self._config.brand_end_slide_path)
			brand_inserted = True

		prs.save(output_path)
		# Scrub relationship parts that python-pptx doesn't expose.
		_scrub_pptx_zip_xml(output_path, self._text_rules, hit_counter=hits)

		warnings = ['WARN_FORBIDDEN_REPLACED'] if hits['replaced'] > 0 else []
		if removed > 0:
			warnings.append('WARN_TAIL_REMOVED')
		if brand_inserted:
			warnings.append('WARN_BRAND_APPENDED')

		return CleanOutput(
			aid=etl_out.aid,
			channel_id=etl_out.channel_id,
			clean_pptx_path=output_path,
			forbidden_keywords=list(self._config.forbidden_keywords),
			warnings=warnings,
		)


def run_clean(
	conn,
	*,
	source_batch_id: str,
	etl_out: EtlOutput,
	workshop: WorkshopB,
) -> CleanOutput:
	started_at = datetime.now(timezone.utc)
	try:
		out = workshop.deep_clean(etl_out)
		# Validate the cleaned pptx is still readable and capture final stats.
		prs = Presentation(out.clean_pptx_path)
		pages_count = len(prs.slides)
		file_size_bytes = out.clean_pptx_path.stat().st_size
		file_size_kb = int(math.ceil(file_size_bytes / 1024))
		upsert_processed_asset(
			conn,
			aid=etl_out.aid,
			source_batch_id=source_batch_id,
			fields={
				'local_pptx_path': out.clean_pptx_path,
				'local_cover_path': etl_out.local_cover_path,
				'pages_count': pages_count,
				'file_size_kb': file_size_kb,
				'file_format': etl_out.file_format,
			},
		)
		status = StageStatus.success
		error_code = None
		error_message = None
		warnings = out.warnings
		artifacts = {
			'clean_pptx_path': str(out.clean_pptx_path),
			'forbidden_keywords': out.forbidden_keywords,
			'pages_count': pages_count,
			'file_size_kb': file_size_kb,
		}
	except Exception as exc:  # noqa: BLE001
		out = None
		status = StageStatus.failed
		error_code = 'CLEAN_FAILED'
		error_message = str(exc)
		warnings = []
		artifacts = {}
	finally:
		finished_at = datetime.now(timezone.utc)

	record_stage(
		conn,
		StageRecord(
			aid=etl_out.aid,
			stage=StageName.B,
			status=status,
			started_at=started_at,
			finished_at=finished_at,
			error_code=error_code,
			error_message=error_message,
			warnings=warnings,
			artifacts=artifacts,
		),
	)

	if status != StageStatus.success or out is None:
		raise ValueError(error_message or 'CLEAN failed')

	return out
