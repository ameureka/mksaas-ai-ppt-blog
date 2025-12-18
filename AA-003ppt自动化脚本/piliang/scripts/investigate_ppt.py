import os
import zipfile
import shutil
from pptx import Presentation
from glob import glob

# 配置路径
DOWNLOAD_DIR = "/Users/ameureka/Desktop/v0-mksaas-analycis-1/v0-mk-saas-tools/pachong/downloads/ppt_moban"
TEMP_DIR = "./temp_investigation"

def investigate_ppt():
    # 1. 找一个 ZIP 文件
    zips = glob(os.path.join(DOWNLOAD_DIR, "*.zip"))
    if not zips:
        print("No zip files found in", DOWNLOAD_DIR)
        return

    target_zip = zips[0] # 随机取第一个
    print(f"=== Investigating: {os.path.basename(target_zip)} ===")

    # 2. 解压
    if os.path.exists(TEMP_DIR):
        shutil.rmtree(TEMP_DIR)
    os.makedirs(TEMP_DIR)
    
    ppt_file = None
    with zipfile.ZipFile(target_zip, 'r') as zip_ref:
        zip_ref.extractall(TEMP_DIR)
        for root, dirs, files in os.walk(TEMP_DIR):
            for file in files:
                if file.endswith((".pptx", ".ppt")):
                    ppt_file = os.path.join(root, file)
                    break
    
    if not ppt_file:
        print("No PPTX file found in zip.")
        return

    print(f"Found PPT: {os.path.basename(ppt_file)}")

    # 3. 解析 PPT
    prs = Presentation(ppt_file)
    
    # A. 元数据
    print("\n--- Core Properties ---")
    print(f"Title: {prs.core_properties.title}")
    print(f"Author: {prs.core_properties.author}")
    print(f"Comments: {prs.core_properties.comments}")
    # print(f"Company: {prs.core_properties.company}") # 很多水印藏在这里

    # B. 页面结构分析
    total_slides = len(prs.slides)
    print(f"\n--- Slide Analysis (Total {total_slides}) ---")
    
    for i, slide in enumerate(prs.slides):
        print(f"\n[Slide {i+1}/{total_slides}]")
        
        # 提取文本
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    texts.append(text[:50].replace('\n', ' ')) # 只显示前50字符
        
        print(f"  Texts ({len(texts)}): {texts}")
        
        # 特殊检查：是否包含敏感词
        sensitive = [t for t in texts if "第一PPT" in t or "1ppt" in t]
        if sensitive:
            print(f"  ⚠️  WARNING: Found sensitive words: {sensitive}")

    # C. 母版分析
    print(f"\n--- Master Analysis ({len(prs.slide_masters)}) Masters ---")
    for i, master in enumerate(prs.slide_masters):
        print(f"[Master {i+1}] Layouts: {len(master.slide_layouts)}")
        # 检查母版里的文字
        master_texts = []
        for shape in master.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                master_texts.append(shape.text.strip())
        print(f"  Master Texts: {master_texts}")

if __name__ == "__main__":
    investigate_ppt()
