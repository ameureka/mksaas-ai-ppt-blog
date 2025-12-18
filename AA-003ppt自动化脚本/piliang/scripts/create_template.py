from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_brand_template():
    prs = Presentation()
    
    # 使用空白版式 (通常是索引 6，但也可能不同，我们清除所有占位符自己画)
    slide_layout = prs.slide_layouts[6] 
    slide = prs.slides.add_slide(slide_layout)

    # 1. 设置背景色 (深色背景显高端，或者白色背景通用)
    # 这里我们用白色背景，文字用深色，比较百搭
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    # 2. 大标题: THANKS
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(8)
    height = Inches(1.5)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "THANKS"
    
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.bold = True
    p.font.size = Pt(80)
    p.font.color.rgb = RGBColor(50, 50, 50) # 深灰
    p.font.name = 'Arial'

    # 3. 副标题: 感谢观看
    top = Inches(4)
    height = Inches(1)
    txBox2 = slide.shapes.add_textbox(left, top, width, height)
    tf2 = txBox2.text_frame
    tf2.text = "感谢您的观看"
    
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(32)
    p2.font.color.rgb = RGBColor(100, 100, 100) # 中灰
    p2.font.name = 'Microsoft YaHei' # 微软雅黑

    # 4. 品牌链接: PPTHub
    top = Inches(6)
    height = Inches(1)
    txBox3 = slide.shapes.add_textbox(left, top, width, height)
    tf3 = txBox3.text_frame
    tf3.text = "更多精品模板: https://www.ppthub.shop/"
    
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    p3.font.size = Pt(20)
    p3.font.color.rgb = RGBColor(0, 112, 192) # 品牌蓝 (示例)
    p3.font.name = 'Arial'
    
    # 保存
    project_root = Path(__file__).resolve().parents[1]
    save_path = project_root / "templates" / "brand_end_slide.pptx"
    save_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(save_path))
    print(f"Brand template saved to {save_path}")

if __name__ == "__main__":
    create_brand_template()
